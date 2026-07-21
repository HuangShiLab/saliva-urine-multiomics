# -*- coding: utf-8 -*-
"""Build GTDB vs HROM database-comparison PPTX (CJK)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

ROOT = "/Users/zhangyf/Projects/baiyunyang"
OUT  = os.path.join(ROOT, "Results", "GTDB_vs_HROM_comparison.pptx")
FONT = "PingFang SC"

# ---- palette ----
DARK  = RGBColor(0x0B,0x3B,0x47)   # deep teal (title/closing)
GTDB  = RGBColor(0x2E,0x6F,0x95)   # steel blue
HROM  = RGBColor(0x1B,0x99,0x8B)   # teal-green
ACC   = RGBColor(0xC0,0x39,0x2B)   # pathogen red
INK   = RGBColor(0x20,0x2A,0x30)   # body text
MUTE  = RGBColor(0x6B,0x7A,0x82)   # muted
WHITE = RGBColor(0xFF,0xFF,0xFF)
PANEL = RGBColor(0xEE,0xF4,0xF5)   # light teal-gray
PANELB= RGBColor(0xE2,0xEC,0xEE)   # slightly darker panel
GOLD  = RGBColor(0xD9,0xA4,0x41)

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    # send to back
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s

def _set_font(run, size, bold, color, font=FONT, italic=False):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.name = font; run.font.color.rgb = color
    rPr = run.font._rPr
    for tag in ("a:ea","a:cs"):
        e = rPr.makeelement(qn(tag), {"typeface": font}); rPr.append(e)

def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {}); el.append(ef)
        sh = ef.makeelement(qn('a:outerShdw'), {'blurRad':'40000','dist':'20000','dir':'5400000','rotWithShape':'0'}); ef.append(sh)
        c = sh.makeelement(qn('a:srgbClr'), {'val':'1F2933'}); sh.append(c)
        a = c.makeelement(qn('a:alpha'), {'val':'24000'}); c.append(a)
    return sp

def text(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         pad=0.0, line_spacing=1.0):
    """lines: list of (txt,size,bold,color[,italic]) OR list of paragraphs where
    each paragraph is a list of such run-tuples."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf, m, Inches(pad))
    first = True
    for para in lines:
        runs = para if isinstance(para, list) else [para]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.line_spacing = line_spacing
        p.space_after = Pt(2); p.space_before = Pt(0)
        for rt in runs:
            txt, size, bold, color = rt[0], rt[1], rt[2], rt[3]
            italic = rt[4] if len(rt) > 4 else False
            r = p.add_run(); r.text = txt; _set_font(r, size, bold, color, italic=italic)
    return tb

def chip(s, x, y, w, h, label, fill, fg=WHITE, size=12, bold=True):
    sp = box(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try: sp.adjustments[0] = 0.5
    except Exception: pass
    tf = sp.text_frame; tf.word_wrap = True
    for m in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf, m, Inches(0.03))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; _set_font(r, size, bold, fg)
    sp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sp

def header(s, title, kicker=None, num=None):
    box(s, 0, 0, SW, 0.16, fill=HROM)            # slim top motif bar
    box(s, 0, 0.16, 0.20, 1.0, fill=GTDB)        # small left tab accent
    ty = 0.42
    if kicker:
        text(s, 0.55, 0.30, 11.5, 0.3, [[(kicker, 12, True, HROM)]])
        ty = 0.58
    text(s, 0.55, ty, 12.2, 0.9, [[(title, 27, True, DARK)]])
    if num is not None:
        text(s, 12.4, 7.02, 0.8, 0.3, [[(str(num), 10, False, MUTE)]], align=PP_ALIGN.RIGHT)

def img_fit(s, path, x, y, max_w, max_h, center_x=None):
    w,h = Image.open(path).size; ar = h/w
    iw = max_w; ih = iw*ar
    if ih > max_h: ih = max_h; iw = ih/ar
    if center_x is not None: x = center_x - iw/2
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(iw), Inches(ih))
    return iw, ih

def table(s, x, y, w, col_w, rows, header_fill=DARK, header_fg=WHITE,
          fsize=12, hsize=12, row_h=0.34, header_h=0.4, zebra=(WHITE,PANEL),
          align=PP_ALIGN.LEFT, col_colors=None):
    nr=len(rows); nc=len(rows[0])
    gh=[Inches(header_h)]+[Inches(row_h)]*(nr-1)
    tb=s.shapes.add_table(nr,nc,Inches(x),Inches(y),Inches(w),sum(gh,Emu(0))).table
    tb.first_row=False; tb.horz_banding=False
    for j,cw in enumerate(col_w): tb.columns[j].width=Inches(cw)
    for i,row in enumerate(rows):
        tb.rows[i].height=gh[i]
        for j,val in enumerate(row):
            c=tb.cell(i,j); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Inches(0.08); c.margin_right=Inches(0.06)
            c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
            if i==0: c.fill.solid(); c.fill.fore_color.rgb=header_fill
            else: c.fill.solid(); c.fill.fore_color.rgb=zebra[i%2]
            tfc=c.text_frame; tfc.word_wrap=True
            p=tfc.paragraphs[0]; p.alignment = (align if j>0 else PP_ALIGN.LEFT)
            r=p.add_run(); r.text=str(val)
            if i==0: _set_font(r,hsize,True,header_fg)
            else:
                col=INK
                if col_colors and j in col_colors: col=col_colors[j]
                _set_font(r,fsize, j==0, col)
    return tb

def takeaway(s, x, y, w, txt, color=GTDB, h=0.7):
    box(s, x, y, 0.09, h, fill=color)
    text(s, x+0.22, y, w-0.3, h, [[("结论  ", 12, True, color), (txt, 13, False, INK)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

# =========================================================================
# S1 — Title
# =========================================================================
s = slide(DARK)
box(s, 0, 0, SW, 0.18, fill=HROM)
box(s, 0, 6.95, SW, 0.55, fill=RGBColor(0x07,0x2B,0x34))
text(s, 0.8, 1.35, 11.7, 0.4, [[("微生物组数据库对比  ·  Database Benchmark", 14, True, RGBColor(0x7FD3C7,0x00,0x00) if False else RGBColor(0x6F,0xD0,0xC4))]])
text(s, 0.8, 1.95, 11.7, 2.2, [
    [("唾液与尿液微生物组", 40, True, WHITE)],
    [("GTDB ", 40, True, RGBColor(0x8F,0xC1,0xDF)), ("vs ", 40, True, WHITE), ("HROM ", 40, True, RGBColor(0x7F,0xE3,0xD6)), ("注释结果对比", 40, True, WHITE)],
], line_spacing=1.05)
text(s, 0.82, 4.35, 11.5, 0.5, [[("牙周炎–糖尿病–肾病共病队列  ·  同一分析流程  ·  两套参考数据库", 16, False, RGBColor(0xCF,0xE3,0xE6))]])
# stat chips
cy=5.25
for i,(t1,t2,col) in enumerate([("101","受试者 × 唾液+尿液", HROM),("5,186 / 2,931","GTDB / HROM 物种(全表)", GTDB),("2","数据库 · 同一流程", GOLD)]):
    x=0.82+i*3.95
    box(s, x, cy, 3.7, 1.1, fill=RGBColor(0x10,0x47,0x53), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+0.25, cy+0.16, 3.3, 0.55, [[(t1, 26, True, col)]])
    text(s, x+0.27, cy+0.70, 3.25, 0.35, [[(t2, 12, False, RGBColor(0xC6,0xDB,0xDE))]])
text(s, 0.8, 7.04, 11.7, 0.4, [[("同一批样本、同一统计流程，仅更换参考数据库  ·  2026-06-05", 11, False, RGBColor(0x9F,0xBA,0xBE))]])

# =========================================================================
# S2 — Design & purpose
# =========================================================================
s = slide()
header(s, "研究设计与目的", kicker="为什么做两库对比", num=2)
# left: purpose
lx=0.55; ly=1.45; lw=6.3
box(s, lx, ly, lw, 4.9, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, lx+0.3, ly+0.25, lw-0.6, 0.4, [[("目的", 16, True, GTDB)]])
purpose=[
 ("同一批样本、同一套统计流程，仅更换参考数据库 — 分离“数据库效应”。", ),
 ("α/β 多样性 · DS-FDR 两两比较 · Maaslin2 · ANCOM-BC · LEfSe，流程完全一致。", ),
 ("比较两库各自的优缺点与可得结论。", ),
 ("重点:两库识别的差异 marker 及其生物学含义(文献支撑)。", ),
]
yy=ly+0.8
for (t,) in purpose:
    box(s, lx+0.32, yy+0.06, 0.12, 0.12, fill=HROM, shape=MSO_SHAPE.OVAL)
    text(s, lx+0.6, yy-0.04, lw-0.95, 0.7, [[(t, 13.5, False, INK)]], line_spacing=1.05)
    yy+=0.78
# right: design
rx=7.15; rw=5.65
box(s, rx, 1.45, rw, 4.9, fill=WHITE, line=PANELB, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+0.3, 1.68, rw-0.6, 0.4, [[("分组 · 疾病负担梯度", 16, True, GTDB)]])
grad=[("N",34,"正常",RGBColor(0x8C,0x8C,0x8C)),("P",12,"单纯牙周炎",GTDB),
      ("PD",6,"+糖尿病",ACC),("PC",25,"+肾病",HROM),("PCD",24,"+糖尿病肾病",RGBColor(0x6A,0x3D,0x9A))]
gy=2.25
for i,(g,n,d,col) in enumerate(grad):
    x=rx+0.32+i*1.05
    box(s, x, gy, 0.92, 0.92, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, gy+0.14, 0.92, 0.4, [[(g,17,True,WHITE)]], align=PP_ALIGN.CENTER)
    text(s, x, gy+0.5, 0.92, 0.3, [[("n=%d"%n,11,True,WHITE)]], align=PP_ALIGN.CENTER)
    text(s, x-0.05, gy+0.95, 1.02, 0.5, [[(d,9.5,False,MUTE)]], align=PP_ALIGN.CENTER)
text(s, rx+0.3, gy+1.55, rw-0.6, 1.6, [
 [("负担梯度  ", 13, True, GTDB), ("N(0) → P(1) → PD/PC(2) → PCD(3)", 13, False, INK)],
 [("位点  ", 13, True, HROM), ("唾液 S 与尿液 U 各 101 样本(配对)", 13, False, INK)],
 [("派生指示变量  ", 13, True, INK), ("Periodontitis / Diabetes / KidneyDisease (0/1)", 12.5, False, INK)],
 [("数据性质  ", 13, True, INK), ("两库均为相对丰度(列和=1)", 12.5, False, INK)],
], line_spacing=1.15)
takeaway(s, 0.55, 6.55, 12.3,
         "唯一变量是参考数据库 → 两库结果的异同可直接归因于注释库本身。", color=HROM, h=0.6)

# =========================================================================
# S3 — Two databases: essence
# =========================================================================
s = slide()
header(s, "两个数据库的本质区别", kicker="GTDB vs HROM", num=3)
rows=[
 ["维度","GTDB","HROM"],
 ["构建依据","全基因组系统发育、rank 标准化\n(Parks 2018/2022)","人类相关 / 口腔参考基因组目录\n(curated)"],
 ["覆盖范围","全细菌 / 古菌域,通用","人体 / 口腔为主"],
 ["命名特点","规范种名 + 大量 spXXXXXX\n占位 / 未培养基因组","命名种 s__ + 大量基因组 bins\ng__…_HROM_Genome_####"],
 ["物种数(全表)","5,186","2,931"],
 ["优势场景","跨研究可比、对接机制文献","口腔样本敏感、低丰度口腔致病菌"],
]
table(s, 0.55, 1.5, 12.25, [2.3,5.0,4.95], rows,
      header_fill=DARK, fsize=12.5, hsize=13, row_h=0.78, header_h=0.42,
      col_colors={1:GTDB,2:HROM})
takeaway(s, 0.55, 6.55, 12.3,
   "GTDB = 通用·标准化·规范命名;HROM = 口腔优化·敏感但多为基因组 bins。二者互补。", color=GTDB, h=0.6)

# =========================================================================
# S4 — Venn / overlap
# =========================================================================
s = slide()
header(s, "物种检出与跨库重叠", kicker="Venn · 同一位点检出物种", num=4)
img_fit(s, os.path.join(ROOT,"Results","venn_saliva.png"), 0.55, 1.5, 4.7, 3.5)
img_fit(s, os.path.join(ROOT,"Results","venn_urine.png"),  5.35, 1.5, 4.7, 3.5)
text(s, 0.55, 4.95, 4.7, 0.3, [[("唾液 Saliva", 13, True, INK)]], align=PP_ALIGN.CENTER)
text(s, 5.35, 4.95, 4.7, 0.3, [[("尿液 Urine", 13, True, INK)]], align=PP_ALIGN.CENTER)
# right numbers panel
px=10.25; pw=2.6
box(s, px, 1.5, pw, 3.75, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, px+0.2, 1.66, pw-0.4, 3.5, [
 [("唾液", 13, True, GTDB)],
 [("共有 ", 12, False, INK), ("795", 16, True, ACC)],
 [("GTDB 特有 ", 11.5, False, MUTE), ("2,662", 13, True, GTDB)],
 [("HROM 特有 ", 11.5, False, MUTE), ("2,121", 13, True, HROM)],
 [("", 6, False, INK)],
 [("尿液", 13, True, GTDB)],
 [("共有 ", 12, False, INK), ("265", 16, True, ACC)],
 [("GTDB 特有 ", 11.5, False, MUTE), ("2,398", 13, True, GTDB)],
 [("HROM 特有 ", 11.5, False, MUTE), ("159", 13, True, HROM)],
], line_spacing=1.12)
takeaway(s, 0.55, 5.55, 12.3,
  "跨库物种名重叠低,主因命名体系不同(GTDB sp-accession vs HROM s__/genome bins),并非生物学不一致 → 比较应落到属 / 功能层面交叉印证。",
  color=HROM, h=1.0)

# =========================================================================
# S5 — Richness
# =========================================================================
s = slide()
header(s, "物种丰富度:位点依赖的“此消彼长”", kicker="每样本检出物种数", num=5)
img_fit(s, os.path.join(ROOT,"Results","richness_boxplot.png"), 0.55, 1.5, 7.7, 4.5)
# right table/callouts
rx=8.5; rw=4.35
box(s, rx, 1.5, rw, 4.4, fill=PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, rx+0.25, 1.7, rw-0.5, 0.4, [[("中位检出物种 / 样本", 14, True, GTDB)]])
table(s, rx+0.22, 2.25, rw-0.44, [1.5,1.2,1.18],
      [["位点","HROM","GTDB"],["唾液","1,195","816"],["尿液","28","121"]],
      header_fill=DARK, fsize=13, hsize=12.5, row_h=0.5, header_h=0.42,
      align=PP_ALIGN.CENTER, col_colors={1:HROM,2:GTDB})
text(s, rx+0.25, 4.0, rw-0.5, 1.7, [
 [("唾液  ", 13, True, HROM), ("HROM ≫ GTDB", 13, True, INK)],
 [("尿液  ", 13, True, GTDB), ("GTDB ≫ HROM", 13, True, INK)],
 [("配对 Wilcoxon ", 11.5, False, MUTE), ("p < 1e-17", 12.5, True, ACC)],
], line_spacing=1.2)
takeaway(s, 0.55, 6.2, 12.3,
  "HROM 口腔优化 → 唾液检出更全;GTDB 通用(含环境/皮肤)→ 尿液(低生物量)检出更全。库的“主场”决定丰富度。",
  color=GTDB, h=0.95)

# =========================================================================
# S6 — Beta / PERMANOVA
# =========================================================================
s = slide()
header(s, "群落结构:两库结论一致", kicker="Beta 多样性 · PERMANOVA", num=6)
# left callouts
def statcard(x,y,w,h,big,small,col):
    box(s,x,y,w,h,fill=WHITE,line=PANELB,line_w=1.2,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s,x,y,0.12,h,fill=col)
    text(s,x+0.3,y+0.12,w-0.45,0.6,[[(big,24,True,col)]])
    text(s,x+0.32,y+0.62,w-0.5,h-0.6,[[(small,12,False,INK)]],line_spacing=1.05)
statcard(0.55,1.55,6.1,1.25,"位点是最大变异来源",
   "合并样品 PERMANOVA — Position R²:  HROM 0.49 (Bray) / 0.74 (rPCA)；GTDB 0.44 / 0.77。", HROM)
statcard(0.55,2.98,6.1,1.25,"牙周炎:群落效应小但方向一致",
   "唾液 Periodontitis 边际 R² ≈ 0.015 (Bray, p≈0.06–0.08),HROM 与 GTDB 几乎一致。", GTDB)
text(s,0.6,4.45,6.1,1.0,[
 [("• ", 13, True, GTDB),("两库在排序(PCoA/rPCA)上呈现相同的“位点两大云团”结构。",12.5,False,INK)],
 [("• ", 13, True, GTDB),("疾病在群落整体水平信号弱 → marker 需靠 per-taxon 差异分析(见后)。",12.5,False,INK)],
], line_spacing=1.15)
# right: two combined PCoA stacked
text(s, 6.9, 1.5, 6.0, 0.28, [[("HROM — 合并 PCoA (Bray | rPCA)", 11, True, HROM)]])
img_fit(s, os.path.join(ROOT,"Results","figures","combined_pcoa.png"), 6.9, 1.85, 6.0, 1.95)
text(s, 6.9, 3.95, 6.0, 0.28, [[("GTDB — 合并 PCoA (Bray | rPCA)", 11, True, GTDB)]])
img_fit(s, os.path.join(ROOT,"GTDB_Results","figures","combined_pcoa.png"), 6.9, 4.3, 6.0, 1.95)
takeaway(s, 0.55, 6.3, 12.3,
  "无论用哪个库,采样位点主导、牙周炎群落效应小而一致 → 群落层面结论稳健、不依赖数据库。", color=HROM, h=0.85)

# =========================================================================
# S7 — Periodontitis signature markers (CORE)
# =========================================================================
s = slide()
header(s, "牙周炎 signature:两库 marker 对比", kicker="N vs P · 唾液 · DS-FDR", num=7)
# sig counts
box(s,0.55,1.45,5.9,0.95,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,0.8,1.55,2.9,0.8,[[("HROM",12,True,HROM)],[("132",30,True,HROM),("  显著",12,False,MUTE)]],line_spacing=1.0)
text(s,3.55,1.55,2.7,0.8,[[("GTDB",12,True,GTDB)],[("79",30,True,GTDB),("  显著",12,False,MUTE)]],line_spacing=1.0)
text(s,0.6,2.42,5.9,0.3,[[("显著差异物种数 (q<0.05;q<0.10:161 / 94)",10.5,False,MUTE)]])
# two columns of markers
hh=2.55
text(s,6.7,1.45,3.0,0.35,[[("HROM Top markers",13,True,HROM)]])
text(s,9.95,1.45,3.0,0.35,[[("GTDB Top markers",13,True,GTDB)]])
hrom_m=["g__Treponema_B (genome bins)","Treponema_D parvum","Porphyromonas gulae  ★","Porphyromonas gingivalis","Treponema_B putidum","CAJPNN01 sp905372035"]
gtdb_m=["Treponema_D parvum","Porphyromonas gingivalis  ★","Treponema_B denticola (红复合体)","Filifactor alocis","Prevotella koreensis","Treponema sp905373565"]
yy=1.85
for a,b in zip(hrom_m,gtdb_m):
    box(s,6.7,yy+0.05,3.05,0.42,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s,6.85,yy+0.06,2.8,0.42,[[(a,11,False,INK)]],anchor=MSO_ANCHOR.MIDDLE)
    box(s,9.95,yy+0.05,2.9,0.42,fill=PANELB,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s,10.1,yy+0.06,2.65,0.42,[[(b,11,False,INK)]],anchor=MSO_ANCHOR.MIDDLE)
    yy+=0.52
# left lower note
text(s,0.55,3.05,5.9,2.4,[
 [("两库共同指向 ",13,True,INK),("红 / 橙复合体牙周致病菌",13,True,ACC)],
 [("",6,False,INK)],
 [("• Treponema(密螺旋体)",12.5,False,INK)],
 [("• Porphyromonas(卟啉单胞菌)",12.5,False,INK)],
 [("• Filifactor alocis / Prevotella",12.5,False,INK)],
 [("",4,False,INK)],
 [("★ = 各库居首的 Porphyromonas 命中,命名不同(见下页)",11,True,ACC)],
], line_spacing=1.18)
takeaway(s, 0.55, 6.4, 12.3,
  "两库高度一致地识别经典牙周致病菌;HROM 在唾液检出更多(口腔优化),GTDB 给出更规范的种名。", color=HROM, h=0.7)

# =========================================================================
# S8 — Naming / resolution difference (EMPHASIS)
# =========================================================================
s = slide()
header(s, "同一信号,不同命名分辨率", kicker="marker 解读的关键差异", num=8)
# two columns
cw=6.05
box(s,0.55,1.5,cw,3.3,fill=WHITE,line=GTDB,line_w=1.6,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
chip(s,0.8,1.72,1.6,0.42,"GTDB",GTDB)
text(s,0.8,2.35,cw-0.5,2.3,[
 [("规范、可培养种名 → 直接对接文献与机制",13,True,GTDB)],
 [("• Porphyromonas gingivalis(角石病原)",12.5,False,INK)],
 [("• Treponema denticola(红复合体)",12.5,False,INK)],
 [("• Filifactor alocis(新兴致病菌)",12.5,False,INK)],
 [("代价:大量 spXXXX 占位、属被拆成 _A/_B 后缀",11.5,False,MUTE)],
], line_spacing=1.22)
box(s,6.75,1.5,cw,3.3,fill=WHITE,line=HROM,line_w=1.6,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
chip(s,7.0,1.72,1.6,0.42,"HROM",HROM)
text(s,7.0,2.35,cw-0.5,2.3,[
 [("口腔特异,但多为基因组 bins",13,True,HROM)],
 [("• g__Treponema_B_HROM_Genome_####(无培养种名)",12,False,INK)],
 [("• 居首的 Porphyromonas = P. gulae(P. gingivalis 紧随)",12.5,True,ACC)],
 [("  (与 P. gingivalis 高度同源、毒力相似;动物相关/人兽共患)",11.5,False,MUTE)],
 [("优点:捕获更多低丰度口腔类群",11.5,False,MUTE)],
], line_spacing=1.22)
# center callout
box(s,0.55,5.05,12.3,1.05,fill=DARK,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,0.85,5.18,11.7,0.8,[
 [("关键  ",14,True,GOLD),("同一生物学信号、不同命名 —— 居首的 ",13,False,WHITE),
  ("Porphyromonas",13,True,GOLD),(":HROM = ",13,False,WHITE),
  ("P. gulae",13,True,RGBColor(0x7F,0xE3,0xD6)),(",GTDB = ",13,False,WHITE),
  ("P. gingivalis",13,True,RGBColor(0x8F,0xC1,0xDF)),("(gulae 在 GTDB 跌至 #58)。",12,False,WHITE)],
 [("报告 marker 时须标注跨库对应关系,避免把“命名差异”误读为“生物学差异”。",12.5,False,RGBColor(0xCF,0xE3,0xE6))],
], line_spacing=1.2)
takeaway(s, 0.55, 6.35, 12.3,
  "GTDB 利于机制与跨研究对接;HROM 利于口腔敏感发现 —— 两者交叉验证最稳妥。", color=GTDB, h=0.7)

# =========================================================================
# S9 — Marker biology + literature
# =========================================================================
s = slide()
header(s, "Marker 的生物学含义与文献支撑", kicker="为什么这些菌重要", num=9)
rows=[
 ["Marker","角色 / 机制","共病关联","文献"],
 ["P. gingivalis","牙周“角石病原”(keystone),红复合体;\n操纵宿主免疫致 dysbiosis","入血激活补体;经 NF-κB/\nNLRP3-铁死亡促 CKD","Socransky 1998;\nHajishengallis 2012"],
 ["T. denticola /\nT. forsythia","红复合体,与牙周袋深度强相关","CKD 患者富集","Socransky 1998"],
 ["Filifactor alocis","新兴牙周致病菌,P. gingivalis 协同","糖尿病者牙周炎严重度↑","Microbes Infect 2015;\nFront Dent Med 2024"],
 ["P. gulae","动物相关 Porphyromonas,\n毒力/免疫特征似 P. gingivalis","人兽共患潜力","Lenzo 2016;\nFront Microbiol 2023"],
 ["Prevotella 属","橙复合体,牙周炎相关","—","Socransky 1998"],
]
table(s, 0.55, 1.45, 12.25, [2.05,4.5,3.0,2.7], rows,
      header_fill=DARK, fsize=11.5, hsize=12, row_h=0.86, header_h=0.4,
      col_colors={0:ACC})
takeaway(s, 0.55, 6.65, 12.3,
  "两库的牙周 marker 均为文献明确的致病菌,且与糖尿病、慢性肾病(CKD)存在机制联系 → 结果生物学可信。",
  color=ACC, h=0.6)

# =========================================================================
# S10 — Maaslin2 multivariable
# =========================================================================
s = slide()
header(s, "多变量校正 marker(adjusted effect)", kicker="Maaslin2 · 唾液 · 控制 Age+Sex 及彼此", num=10)
text(s,0.55,1.45,12.2,0.5,[[("模型  ",12,True,GTDB),("microbe ~ Periodontitis + Diabetes + KidneyDisease + Age + Sex  ",12,False,INK),("(显著:q<0.25)",11,False,MUTE)]])
table(s, 0.55, 2.05, 7.4, [3.0,2.2,2.2],
      [["疾病因子","HROM 显著菌数","GTDB 显著菌数"],
       ["Periodontitis","94","59"],["Diabetes","42","38"],["KidneyDisease","113","70"]],
      header_fill=DARK, fsize=13, hsize=12.5, row_h=0.6, header_h=0.45,
      align=PP_ALIGN.CENTER, col_colors={1:HROM,2:GTDB})
box(s,8.2,2.05,4.6,2.95,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,8.45,2.25,4.15,2.7,[
 [("解读",14,True,GTDB)],
 [("• 控制混杂后,两库均检出牙周/肾病独立关联菌。",12,False,INK)],
 [("• HROM 在唾液检出更多(口腔敏感)。",12,False,INK)],
 [("• KidneyDisease 关联菌最多 → ",12,False,INK),("口腔–系统(肾)轴",12,True,ACC),("线索。",12,False,INK)],
], line_spacing=1.25)
takeaway(s, 0.55, 5.45, 12.3,
  "唾液中“肾病相关”口腔菌数量可观,呼应红复合体在 CKD 富集的文献 → 支持口腔–肾轴假说,值得后续验证。",
  color=GTDB, h=0.95)

# =========================================================================
# S11 — Pros / Cons
# =========================================================================
s = slide()
header(s, "两库优缺点汇总", kicker="如何取舍", num=11)
def proscons(x,title,col,pros,cons):
    w=6.05
    box(s,x,1.5,w,4.8,fill=WHITE,line=col,line_w=1.6,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    chip(s,x+0.25,1.72,1.7,0.45,title,col,size=13)
    text(s,x+0.28,2.4,w-0.55,0.32,[[("优点",12.5,True,col)]])
    yy=2.78
    for t in pros:
        box(s,x+0.32,yy+0.07,0.1,0.1,fill=col,shape=MSO_SHAPE.OVAL)
        text(s,x+0.55,yy-0.03,w-0.85,0.6,[[(t,11.8,False,INK)]],line_spacing=1.05); yy+=0.5
    text(s,x+0.28,yy+0.05,w-0.55,0.32,[[("缺点",12.5,True,ACC)]]); yy+=0.4
    for t in cons:
        box(s,x+0.32,yy+0.07,0.1,0.1,fill=ACC,shape=MSO_SHAPE.OVAL)
        text(s,x+0.55,yy-0.03,w-0.85,0.6,[[(t,11.8,False,INK)]],line_spacing=1.05); yy+=0.5
proscons(0.55,"GTDB",GTDB,
 ["标准化·基因组系统发育·通用","规范种名 → 跨研究可比、对接机制","尿液/低生物量类群覆盖更全"],
 ["大量未培养占位名(spXXXX/属代码)难解读","非口腔优化 → 唾液口腔菌可能低估","种被拆成 _A/_B 后缀,增加解读成本"])
proscons(6.75,"HROM",HROM,
 ["口腔/人体优化 → 唾液敏感","检出更多低丰度口腔致病菌","含口腔特异基因组(MAG/bins)"],
 ["大量基因组 bins 无培养种名","尿液/非口腔覆盖差(尿液丰富度低)","跨库/跨研究可比性差"])
takeaway(s, 0.55, 6.45, 12.3,
  "没有“更好的库”,只有“更合适的库”:唾液选 HROM(敏感),尿液选 GTDB(覆盖),命名靠 GTDB。", color=GTDB, h=0.65)

# =========================================================================
# S12 — Conclusions each DB yields
# =========================================================================
s = slide()
header(s, "各数据库能得到的结论", kicker="共同 · 独到", num=12)
def concard(x,y,w,h,title,col,items,off=0.82,sp=0.56):
    box(s,x,y,w,h,fill=PANEL if col!=DARK else DARK,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s,x,y,w,0.12,fill=col if col!=DARK else GOLD)
    fg = WHITE if col==DARK else INK
    text(s,x+0.28,y+0.24,w-0.55,0.4,[[(title,15,True,(GOLD if col==DARK else col))]])
    yy=y+off
    for t in items:
        text(s,x+0.3,yy,w-0.6,0.55,[[("• ",12,True,(GOLD if col==DARK else col)),(t,12,False,fg)]],line_spacing=1.05); yy+=sp
concard(0.55,1.5,12.3,2.45,"共同结论(两库一致 → 可靠,交叉验证)",DARK,[
 "采样位点(唾液 vs 尿液)是群落变异的最大来源。",
 "牙周炎 signature = 红/橙复合体(P. gingivalis · T. denticola · Treponema · Prevotella · F. alocis)。",
 "尿液信号弱而稀疏;PD(n=6)及双重共病的附加效应弱、需谨慎。",
])
concard(0.55,4.05,6.05,2.35,"GTDB 独到",GTDB,[
 "以规范种名直接命名 red-complex 与 F. alocis。",
 "便于机制解读与跨研究/跨库对接。",
 "尿液类群覆盖更全。",
],off=0.78,sp=0.52)
concard(6.8,4.05,6.05,2.35,"HROM 独到",HROM,[
 "更高灵敏度:唾液检出更多牙周致病菌。",
 "捕获口腔特异基因组(含 P. gulae 等 bins)。",
 "更适合口腔疾病 marker 发现。",
],off=0.78,sp=0.52)

# =========================================================================
# S13 — Recommendations
# =========================================================================
s = slide()
header(s, "建议:如何用好两个库", kicker="落地策略", num=13)
recs=[
 ("唾液(口腔)","HROM 做敏感发现 + GTDB 做规范命名/可比,互补使用。",HROM),
 ("尿液","以 GTDB 为主(覆盖更全、低生物量更稳)。",GTDB),
 ("marker 报告","标注跨库命名对应(P. gulae ↔ P. gingivalis;Treponema_B/_D ↔ T. denticola 等)。",ACC),
 ("下游分析","两库交叉验证一致的 marker 优先纳入预测模型 / 机制研究。",GOLD),
]
yy=1.6
for i,(h,t,col) in enumerate(recs):
    box(s,0.55,yy,12.3,1.12,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s,0.55,yy,0.12,1.12,fill=col)
    box(s,0.85,yy+0.26,0.6,0.6,fill=col,shape=MSO_SHAPE.OVAL)
    text(s,0.85,yy+0.3,0.6,0.5,[[(str(i+1),20,True,WHITE)]],align=PP_ALIGN.CENTER)
    text(s,1.7,yy+0.16,2.7,0.8,[[(h,15,True,col)]],anchor=MSO_ANCHOR.MIDDLE)
    text(s,4.4,yy+0.16,8.2,0.85,[[(t,13.5,False,INK)]],anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.1)
    yy+=1.25

# =========================================================================
# S14 — Methods & limitations
# =========================================================================
s = slide()
header(s, "方法与局限", kicker="可重复性 · 审稿要点", num=14)
box(s,0.55,1.5,6.05,4.5,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,0.85,1.72,5.5,4.2,[
 [("方法",15,True,GTDB)],
 [("• 相对丰度;≥10% 流行率过滤。",12.5,False,INK)],
 [("• α 多样性(vegan):Observed/Shannon/Simpson/Pielou。",12.5,False,INK)],
 [("• β:Bray-Curtis 与 robust Aitchison(rPCA);PERMANOVA。",12.5,False,INK)],
 [("• 差异丰度:两两 Mann-Whitney + DS-FDR(Jiang 2017)。",12.5,False,INK)],
 [("• 多变量:Maaslin2 + ANCOM-BC(BH);LEfSe(LDA>2)。",12.5,False,INK)],
 [("• 两库流程完全一致,仅注释库不同。",12.5,True,GTDB)],
], line_spacing=1.25)
box(s,6.8,1.5,6.05,4.5,fill=WHITE,line=ACC,line_w=1.4,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,7.1,1.72,5.5,4.2,[
 [("局限",15,True,ACC)],
 [("• 跨库物种名仅按字符串匹配 → 低估真实重叠(应在属/功能层面比较)。",12.5,False,INK)],
 [("• 数据为相对丰度,非原始计数。",12.5,False,INK)],
 [("• PD 组 n=6,为最大统计瓶颈。",12.5,False,INK)],
 [("• 仅 Age/Sex 可校正(无 BMI/吸烟/eGFR/HbA1c)。",12.5,False,INK)],
 [("• 除 P 外各组伴其他慢病(Ch)→ 反映“共病状态”,非纯净单病效应。",12.5,False,INK)],
], line_spacing=1.28)
takeaway(s, 0.55, 6.25, 12.3,
  "结论应表述为 disease-associated comorbidity states;marker 的跨库一致性比单库显著性更可信。",
  color=ACC, h=0.95)

# =========================================================================
# S15 — References
# =========================================================================
s = slide(DARK)
box(s, 0, 0, SW, 0.16, fill=HROM)
text(s, 0.7, 0.5, 11.0, 0.7, [[("参考文献", 26, True, WHITE)]])
refs=[
 "Socransky SS, et al. Microbial complexes in subgingival plaque. J Clin Periodontol. 1998;25(2):134–144.",
 "Hajishengallis G, Darveau RP, Curtis MA. The keystone-pathogen hypothesis. Nat Rev Microbiol. 2012;10(10):717–725.",
 "Filifactor alocis – a new emerging periodontal pathogen. Microbes Infect. 2015;17(7):517–530.",
 "Filifactor alocis & periodontitis severity in diabetic/non-diabetic patients (narrative review). Front Dent Med. 2024.",
 "Lenzo JC, et al. Porphyromonas gulae has virulence/immunological characteristics similar to P. gingivalis. Infect Immun. 2016;84(9):2575–2585.",
 "Parks DH, et al. A standardized bacterial taxonomy based on genome phylogeny. Nat Biotechnol. 2018;36:996–1004.",
 "Parks DH, et al. GTDB: an ongoing census of bacterial and archaeal diversity. Nucleic Acids Res. 2022;50:D785–D794.",
 "Jiang L, et al. Discrete false-discovery rate improves identification of differentially abundant microbes. mSystems. 2017;2(6):e00092-17.",
 "Periodontitis & chronic kidney disease: epidemiology and biological mechanisms (reviews). Front Microbiol. 2021; J Clin Med. 2023.",
]
yy=1.5
for i,r in enumerate(refs):
    text(s,0.7,yy,12.0,0.6,[[("%d. "%(i+1),12,True,RGBColor(0x7F,0xE3,0xD6)),(r,12,False,RGBColor(0xDF,0xEA,0xEC))]],line_spacing=1.05)
    yy+=0.58
text(s,0.7,7.0,12,0.4,[[("数据库:GTDB(genome-based, universal) · HROM(human/oral curated)  |  报告:hrom_group_comparison.html · GTDB_group_comparison.html",10.5,False,RGBColor(0x9F,0xBA,0xBE))]])

prs.save(OUT)
print("SAVED", OUT, "slides=", len(prs.slides._sldIdLst))
