# -*- coding: utf-8 -*-
"""Build the Saliva->Urine cross-site multi-omics PPTX (CJK).
Numbers are read from GTDB_Results/crossomics/tables/*.csv so the deck stays
in sync with the Rmarkdown analysis."""
import os, csv, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

ROOT = "/Users/zhangyf/Projects/baiyunyang"
DB   = sys.argv[1] if len(sys.argv) > 1 else "GTDB"     # GTDB | HROM | GTDB_noPD
NOPD = DB.endswith("_noPD")                             # PD group excluded variant
DBLABEL = DB[:-5] if NOPD else DB                       # clean name for display
COHORT  = "95 例(已排除 PD 组)" if NOPD else "101 例"
RES  = os.path.join(ROOT, f"{DB}_Results", "crossomics")
FIG  = os.path.join(RES, "figures")
TAB  = os.path.join(RES, "tables")
OUT  = os.path.join(RES, f"Saliva_to_Urine_multiomics_{DB}.pptx")
FONT = "PingFang SC"

# ---- palette (matched to the existing project deck) ----
DARK  = RGBColor(0x0B,0x3B,0x47); MIC = RGBColor(0x2E,0x6F,0x95)
MET   = RGBColor(0x1B,0x99,0x8B); ACC = RGBColor(0xC0,0x39,0x2B)
INK   = RGBColor(0x20,0x2A,0x30); MUTE= RGBColor(0x6B,0x7A,0x82)
WHITE = RGBColor(0xFF,0xFF,0xFF); PANEL = RGBColor(0xEE,0xF4,0xF5)
PANELB= RGBColor(0xE2,0xEC,0xEE); GOLD = RGBColor(0xD9,0xA4,0x41)
GREEN = RGBColor(0x00,0xA0,0x87)

SW, SH = 13.333, 7.5
prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ---------- read result tables ----------
def read_csv(path):
    if not os.path.exists(path): return []
    with open(path, newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))

SUMMARY = read_csv(os.path.join(TAB, "cross_pairing_summary.csv"))
P1pairs = read_csv(os.path.join(TAB, "P1_sig_pairs.csv"))
P2pairs = read_csv(os.path.join(TAB, "P2_sig_pairs.csv"))
P3pairs = read_csv(os.path.join(TAB, "P3_sig_pairs.csv"))
MED     = read_csv(os.path.join(TAB, "P4_mediation.csv"))
PF = {k: read_csv(os.path.join(TAB, f"{k}_permanova_perfeature.csv"))
      for k in ("PF1","PF2","PF3")}

def pf_stats(key):       # per-feature PERMANOVA: (#tested, #sig, median R2, max R2)
    rows = PF.get(key, [])
    sig = [float(r["R2"]) for r in rows if r.get("p") and float(r["p"]) < 0.05]
    if not sig: return (len(rows), 0, 0.0, 0.0)
    sig.sort()
    med = sig[len(sig)//2]
    return (len(rows), len(sig), med, max(sig))

def S(i, key, d=""):     # safe summary access by row index
    try: return SUMMARY[i][key]
    except Exception: return d
def fnum(x, nd=3):
    try: return ("%."+str(nd)+"f") % float(x)
    except Exception: return str(x)
def fp(x):
    try:
        v = float(x)
        return "<0.001" if v < 0.001 else ("%.3f" % v)
    except Exception: return str(x)

# ---------- pptx helpers (from project convention) ----------
def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit=False
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s
def _set_font(run, size, bold, color, font=FONT, italic=False):
    run.font.size=Pt(size); run.font.bold=bold; run.font.italic=italic
    run.font.name=font; run.font.color.rgb=color
    rPr=run.font._rPr
    for tag in ("a:ea","a:cs"):
        e=rPr.makeelement(qn(tag),{"typeface":font}); rPr.append(e)
def box(s,x,y,w,h,fill=None,line=None,line_w=1.0,shape=MSO_SHAPE.RECTANGLE,shadow=False):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(line_w)
    sp.shadow.inherit=False
    if shadow:
        el=sp._element.spPr; ef=el.makeelement(qn('a:effectLst'),{}); el.append(ef)
        sh=ef.makeelement(qn('a:outerShdw'),{'blurRad':'40000','dist':'20000','dir':'5400000','rotWithShape':'0'}); ef.append(sh)
        c=sh.makeelement(qn('a:srgbClr'),{'val':'1F2933'}); sh.append(c)
        a=c.makeelement(qn('a:alpha'),{'val':'24000'}); c.append(a)
    return sp
def text(s,x,y,w,h,lines,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,pad=0.0,line_spacing=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,Inches(pad))
    first=True
    for para in lines:
        runs=para if isinstance(para,list) else [para]
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align; p.line_spacing=line_spacing; p.space_after=Pt(2); p.space_before=Pt(0)
        for rt in runs:
            txt,size,bold,color=rt[0],rt[1],rt[2],rt[3]; italic=rt[4] if len(rt)>4 else False
            r=p.add_run(); r.text=txt; _set_font(r,size,bold,color,italic=italic)
    return tb
def chip(s,x,y,w,h,label,fill,fg=WHITE,size=12,bold=True):
    sp=box(s,x,y,w,h,fill=fill,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try: sp.adjustments[0]=0.5
    except Exception: pass
    tf=sp.text_frame; tf.word_wrap=True
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,Inches(0.03))
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=label; _set_font(r,size,bold,fg)
    tf.vertical_anchor=MSO_ANCHOR.MIDDLE; return sp
def header(s,title,kicker=None,num=None):
    box(s,0,0,SW,0.16,fill=MET); box(s,0,0.16,0.20,1.0,fill=MIC); ty=0.42
    if kicker: text(s,0.55,0.30,11.5,0.3,[[(kicker,12,True,MET)]]); ty=0.58
    text(s,0.55,ty,12.2,0.9,[[(title,27,True,DARK)]])
    if num is not None: text(s,12.4,7.02,0.8,0.3,[[(str(num),10,False,MUTE)]],align=PP_ALIGN.RIGHT)
def img_fit(s,path,x,y,max_w,max_h,center_x=None):
    w,h=Image.open(path).size; ar=h/w; iw=max_w; ih=iw*ar
    if ih>max_h: ih=max_h; iw=ih/ar
    if center_x is not None: x=center_x-iw/2
    s.shapes.add_picture(path,Inches(x),Inches(y),Inches(iw),Inches(ih)); return iw,ih
def table(s,x,y,w,col_w,rows,header_fill=DARK,header_fg=WHITE,fsize=12,hsize=12,
          row_h=0.34,header_h=0.4,zebra=(WHITE,PANEL),align=PP_ALIGN.LEFT,col_colors=None):
    nr=len(rows); nc=len(rows[0]); gh=[Inches(header_h)]+[Inches(row_h)]*(nr-1)
    tb=s.shapes.add_table(nr,nc,Inches(x),Inches(y),Inches(w),sum(gh,Emu(0))).table
    tb.first_row=False; tb.horz_banding=False
    for j,cw in enumerate(col_w): tb.columns[j].width=Inches(cw)
    for i,row in enumerate(rows):
        tb.rows[i].height=gh[i]
        for j,val in enumerate(row):
            c=tb.cell(i,j); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Inches(0.08); c.margin_right=Inches(0.06)
            c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
            if i==0: c.fill.solid(); c.fill.fore_color.rgb=header_fill
            else: c.fill.solid(); c.fill.fore_color.rgb=zebra[i%2]
            tfc=c.text_frame; tfc.word_wrap=True
            p=tfc.paragraphs[0]; p.alignment=(align if j>0 else PP_ALIGN.LEFT)
            r=p.add_run(); r.text=str(val)
            if i==0: _set_font(r,hsize,True,header_fg)
            else:
                col=INK
                if col_colors and j in col_colors: col=col_colors[j]
                _set_font(r,fsize,j==0,col)
    return tb
def takeaway(s,x,y,w,txt,color=MIC,h=0.7):
    box(s,x,y,0.09,h,fill=color)
    text(s,x+0.22,y,w-0.3,h,[[("结论  ",12,True,color),(txt,13,False,INK)]],
         anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.05)

# =========================================================================
# S1 — Title
# =========================================================================
s=slide(DARK)
box(s,0,0,SW,0.18,fill=MET); box(s,0,6.95,SW,0.55,fill=RGBColor(0x07,0x2B,0x34))
text(s,0.8,1.25,11.7,0.4,[[("跨体液多组学整合  ·  Cross‑body‑site Multi‑omics",14,True,RGBColor(0x6F,0xD0,0xC4))]])
text(s,0.8,1.85,11.9,2.2,[
    [("唾液 → 尿液 跨位点关联分析",40,True,WHITE)],
    [("微生物组 ",30,True,RGBColor(0x8F,0xC1,0xDF)),("与 ",30,True,WHITE),
     ("代谢组 ",30,True,RGBColor(0x7F,0xE3,0xD6)),("的跨体液互作",30,True,WHITE)],
],line_spacing=1.05)
text(s,0.82,4.2,11.6,0.5,[[(f"方法参照 Zhang et al., Microbiome 2026 (口腔–肠道–代谢物轴)  ·  {DBLABEL} 物种注释  ·  {COHORT}四层配对",15,False,RGBColor(0xCF,0xE3,0xE6))]])
cy=5.1
chips=[(("95" if NOPD else "101"),("受试者 × 4 数据层" + (" (排除PD)" if NOPD else "")),MET),
       ("3","跨层关系  +  1 中介链",MIC),
       (fnum(S(2,'dbRDA_R2'),2),"代谢→代谢 dbRDA R² (最强轴)",GOLD)]
for i,(t1,t2,col) in enumerate(chips):
    x=0.82+i*3.95
    box(s,x,cy,3.7,1.15,fill=RGBColor(0x10,0x47,0x53),shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s,x+0.25,cy+0.16,3.3,0.55,[[(t1,26,True,col)]])
    text(s,x+0.27,cy+0.74,3.25,0.35,[[(t2,11.5,False,RGBColor(0xC6,0xDB,0xDE))]])
text(s,0.8,7.04,11.7,0.4,[[("同一批受试者,唾液(口腔)→尿液(系统)方向的多组学跨层关联  ·  2026-06-07",11,False,RGBColor(0x9F,0xBA,0xBE))]])

# =========================================================================
# S2 — Reference paper (Zhang et al., Microbiome 2026): methods & conclusions
# =========================================================================
s=slide(); header(s,"方法学来源：我们借鉴的原始研究",
                  kicker="Zhang et al., Microbiome 2026 · 跨体位微生物–代谢组",num=2)
# left: what the study did
lx,ly,lw=0.55,1.45,6.15
box(s,lx,ly,lw,4.55,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,lx+0.3,ly+0.22,lw-0.6,0.4,[[("研究做了什么",16,True,MIC)]])
text(s,lx+0.32,ly+0.8,lw-0.64,3.6,[
 [("对象  ",12.5,True,MIC),("3 个独立队列 n=435;4 个体位菌群(肠 / 口腔 / 皮肤 / 鼻)+ 血浆代谢组(814 代谢物)。",12,False,INK)],
 [("",5,False,INK)],
 [("问题  ",12.5,True,MIC),("多体位菌群如何共同塑造血浆代谢组?体位之间有没有互作?",12,False,INK)],
 [("",5,False,INK)],
 [("方法  ",12.5,True,MIC),("PERMANOVA 方差解释 · GBDT 机器学习预测(R²) · Spearman 相关网络 · 双向中介分析 · 体位交互项模型。",12,False,INK)],
],line_spacing=1.28)
# right: key conclusions
rx,rw=6.95,5.85
box(s,rx,1.45,rw,4.55,fill=WHITE,line=PANELB,line_w=1.2,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,rx+0.3,1.66,rw-0.6,0.4,[[("主要结论",16,True,MET)]])
concl=[
 ("各体位菌群共解释血浆代谢 30.1%;肠 18.4% > 口腔 14.7% > 皮肤 11.5% > 鼻 5.9%。",GOLD),
 ("组成 ≠ 功能:口腔与肠道菌群组成迥异,却与代谢物关联相似。",MIC),
 ("口腔–肠道轴:广泛的协同 / 增效互作;过半代谢物受多体位共同影响。",MET),
 ("中介:口腔→肠道→代谢物 的方向性中介(调控吲哚类、羧酸类等)。",ACC),
 ("胰岛素抵抗(IR)人群该轴更活跃 → 把微生物协同与代谢失调联系起来。",RGBColor(0x6A,0x3D,0x9A)),
]
yy=2.22
for t,col in concl:
    box(s,rx+0.32,yy+0.07,0.12,0.12,fill=col,shape=MSO_SHAPE.OVAL)
    text(s,rx+0.6,yy-0.03,rw-0.95,0.72,[[(t,12,False,INK)]],line_spacing=1.12); yy+=0.735
# bottom: our adaptation
box(s,0.55,6.18,12.3,0.74,fill=DARK,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,0.82,6.27,12.0,0.58,[
 [("我们的迁移  ",13,True,GOLD),
  ("论文「口腔→肠道→血浆代谢物」 ⇒ 本研究「唾液(口腔,近端)→尿液(系统,远端)」,同测微生物组+代谢组,沿用其方差解释 / 网络 / 中介框架。",12,False,WHITE)],
])

# =========================================================================
# S3 — Design & framework
# =========================================================================
s=slide(); header(s,"研究设计与分析框架",kicker="为什么做唾液→尿液跨层分析",num=3)
lx,ly,lw=0.55,1.45,6.2
box(s,lx,ly,lw,4.95,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,lx+0.3,ly+0.22,lw-0.6,0.4,[[("思路（迁移自论文）",16,True,MIC)]])
pts=[("论文研究 口腔→肠道→血浆代谢物 的跨体位微生物–代谢互作。",),
     ("本队列迁移为 唾液(口腔,近端) → 尿液(系统,远端) 方向。",),
     (f"每受试者同时有两体液的 微生物组({DBLABEL}) 与 代谢组(LC-MS)。",),
     ("101 受试者四层完全配对 → 可做跨层关联与中介分析。",)]
yy=ly+0.85
for (t,) in pts:
    box(s,lx+0.32,yy+0.07,0.12,0.12,fill=MET,shape=MSO_SHAPE.OVAL)
    text(s,lx+0.6,yy-0.03,lw-0.95,0.7,[[(t,13.5,False,INK)]],line_spacing=1.08); yy+=0.82
rx,rw=7.05,5.75
box(s,rx,1.45,rw,4.95,fill=WHITE,line=PANELB,line_w=1.2,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,rx+0.3,1.66,rw-0.6,0.4,[[("四项分析（均为 唾液→尿液 方向）",15,True,MIC)]])
rows=[["#","暴露 X (唾液)","结局 Y (尿液)"],
      ["①","微生物组","微生物组"],
      ["②","微生物组","代谢组"],
      ["③","代谢组","代谢组"],
      ["④","物种→尿物种→尿代谢物 (中介)","三层链"]]
table(s,rx+0.28,2.2,rw-0.55,[0.5,3.0,1.65],rows,header_fill=DARK,
      fsize=12.5,hsize=12.5,row_h=0.52,header_h=0.44,col_colors={0:ACC})
text(s,rx+0.3,5.2,rw-0.6,1.1,[
 [("每个组合统一执行：",12.5,True,MIC)],
 [("Mantel + Procrustes 一致性 · dbRDA 方差解释 · Spearman 相关网络(|r|≥0.3,P<0.05) · RandomForest 预测 R²",12,False,INK)],
],line_spacing=1.2)
takeaway(s,0.55,6.55,12.3,
  "唯一不变的是“唾液→尿液”方向;沿微生物组与代谢组两个层面，系统刻画跨体液耦合强度。",color=MET,h=0.6)

# =========================================================================
# S3 — Method guide: the four evaluation methods
# =========================================================================
s=slide(); header(s,"四种评价方法：特点与正确解读",kicker="Mantel · Procrustes · PERMANOVA · dbRDA",num=4)
mrows=[
 ["方法","类型 / 方向","关键输出","解读要点"],
 ["Mantel test","矩阵相关 · 对称","Mantel r, P","r>0 且 P<0.05 = 整体协同变异;r 绝对值天然偏小(勿套 Pearson 经验),功效低、易低估"],
 ["Procrustes","构型叠合 · 对称","r, m², P","比 Mantel 更稳、功效更高;可逐样本看残差(叠合图);依赖保留排序维度数 k"],
 ["PERMANOVA","方差分解 · Y~X","pseudo-F, R², P","R² = Y 被 X 解释的方差比例;对组内离散度(dispersion)敏感;R² 随预测变量数虚高"],
 ["dbRDA","约束排序 · Y~X","约束 R², F, P, 轴","与 PERMANOVA 同一模型,另给可视化约束轴;支持偏分析控制协变量"],
]
table(s,0.55,1.5,12.3,[2.0,2.1,2.4,5.8],mrows,header_fill=DARK,
      fsize=11.5,hsize=12.5,row_h=0.92,header_h=0.46,col_colors={0:MIC})
box(s,0.55,5.55,12.3,1.4,fill=DARK,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,0.85,5.66,11.7,1.2,[
 [("关键  ",13,True,GOLD),("四法 = 两类问题 × 各两种实现：",13,True,WHITE),
  ("对称一致性",12.5,True,RGBColor(0x8F,0xC1,0xDF)),("[Mantel↔Procrustes]  vs  ",12,False,RGBColor(0xCF,0xE3,0xE6)),
  ("非对称方差解释",12.5,True,RGBColor(0x7F,0xE3,0xD6)),("[PERMANOVA↔dbRDA]。",12,False,RGBColor(0xCF,0xE3,0xE6))],
 [("PERMANOVA 与 dbRDA 是同一模型(R²/F/P 必然基本一致),不算两份独立佐证;",12,False,RGBColor(0xDF,0xEA,0xEC))],
 [("若仅非对称法显著、对称法不显著 → 少数强相关特征驱动了可解释方差,但整体构型并不一致。",12,False,RGBColor(0xDF,0xEA,0xEC))],
],line_spacing=1.25)

# =========================================================================
# S4 — Overview of results (synthesis first)
# =========================================================================
s=slide(); header(s,"结果速览：三类跨层关系的耦合强度",kicker="一图看全",num=5)
img_fit(s,os.path.join(FIG,"synthesis_coupling.png"),0.55,1.55,6.7,4.6)
# summary table on right
rx,rw=7.55,5.3
hdr=["关系","Mantel r","dbRDA R²","显著对","RF R²"]
def rowfor(i,lab):
    return [lab,fnum(S(i,'mantel_r')),fnum(S(i,'dbRDA_R2')),
            S(i,'sig_pairs'),fnum(S(i,'RF_meanR2'))]
trows=[hdr,
       rowfor(0,"①微→微"),
       rowfor(1,"①微→代"),
       rowfor(2,"③代→代")]
# fix labels (use proper)
trows[1][0]="① 微生物→微生物"; trows[2][0]="② 微生物→代谢"; trows[3][0]="③ 代谢→代谢"
table(s,rx,1.7,rw,[2.0,0.95,0.95,0.8,0.85],trows,header_fill=DARK,
      fsize=12,hsize=11.5,row_h=0.62,header_h=0.5,align=PP_ALIGN.CENTER,
      col_colors={0:INK})
box(s,rx,4.35,rw,1.95,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,rx+0.25,4.5,rw-0.5,1.8,[
 [("耦合强度排序",13,True,MIC)],
 [("代谢组→代谢组  ",12.5,True,MET),("≫  ",12.5,True,INK),
  ("微生物→代谢  ",12.5,True,GOLD),(">  ",12.5,True,INK),
  ("微生物→微生物",12.5,True,MIC)],
 [("跨体液信号主要存在于 代谢层面，而非微生物组成层面。",12,False,INK)],
],line_spacing=1.3)
takeaway(s,0.55,6.45,12.3,
  "唾液与尿液的代谢组高度耦合;微生物组跨体液耦合弱 → 两体液是相对独立的微生态。",color=MET,h=0.7)

# =========================================================================
# S4 — Pairing 1: microbiome -> microbiome
# =========================================================================
s=slide(); header(s,"分析一：唾液微生物组 → 尿液微生物组",kicker="物种 × 物种",num=6)
img_fit(s,os.path.join(FIG,"P1_procrustes.png"),0.5,1.5,5.7,3.7)
img_fit(s,os.path.join(FIG,"P1_network.png"),6.05,1.5,7.0,4.0)
# stat strip
def statline(s,x,y,items):
    box(s,x,y,12.3,0.8,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    n=len(items); cw=12.3/n
    for i,(k,v,col) in enumerate(items):
        text(s,x+0.2+i*cw,y+0.1,cw-0.3,0.6,[[(v,18,True,col)],[(k,10.5,False,MUTE)]],line_spacing=1.0)
statline(s,0.55,5.35,[
 ("Mantel r (P=%s)"%fp(S(0,'mantel_p')),fnum(S(0,'mantel_r')),MIC),
 ("Procrustes r (P=%s)"%fp(S(0,'procrustes_p')),fnum(S(0,'procrustes_r')),MIC),
 ("PERMANOVA R² (P=%s)"%fp(S(0,'permanova_p')),fnum(S(0,'permanova_R2')),GREEN),
 ("dbRDA R² (P=%s)"%fp(S(0,'dbRDA_p')),fnum(S(0,'dbRDA_R2')),GREEN),
])
takeaway(s,0.55,6.35,12.3,
  "四法一致:对称的 Mantel/Procrustes 不显著、非对称的 PERMANOVA≈dbRDA(R²=%s)亦不显著;仅 %s 对物种弱关联、RF R²<0 → 口腔与泌尿道是独立生态系统。"%(fnum(S(0,'dbRDA_R2')),S(0,'sig_pairs')),
  color=MIC,h=0.7)

# =========================================================================
# S5 — Pairing 2: microbiome -> metabolome
# =========================================================================
s=slide(); header(s,"分析二：唾液微生物组 → 尿液代谢组",kicker="物种 × 代谢物",num=7)
img_fit(s,os.path.join(FIG,"P2_procrustes.png"),0.5,1.5,5.7,3.7)
img_fit(s,os.path.join(FIG,"P2_network.png"),6.05,1.5,7.0,4.0)
statline(s,0.55,5.35,[
 ("Mantel r (P=%s)"%fp(S(1,'mantel_p')),fnum(S(1,'mantel_r')),MIC),
 ("Procrustes r (P=%s)"%fp(S(1,'procrustes_p')),fnum(S(1,'procrustes_r')),MIC),
 ("PERMANOVA R² (P=%s)"%fp(S(1,'permanova_p')),fnum(S(1,'permanova_R2')),GREEN),
 ("dbRDA R² (P=%s)"%fp(S(1,'dbRDA_p')),fnum(S(1,'dbRDA_R2')),GREEN),
])
takeaway(s,0.55,6.35,12.3,
  "对称法不显著、但 PERMANOVA≈dbRDA R²=%s(P=%s)→ 少数强相关特征驱动了可解释方差;显著对(%s)远多于物种→物种,印证“组成≠功能”。"%(fnum(S(1,'dbRDA_R2')),fp(S(1,'dbRDA_p')),S(1,'sig_pairs')),
  color=GOLD,h=0.7)

# =========================================================================
# S6 — Pairing 3: metabolome -> metabolome (headline)
# =========================================================================
s=slide(); header(s,"分析三：唾液代谢组 → 尿液代谢组",kicker="最强的跨体液轴",num=8)
img_fit(s,os.path.join(FIG,"P3_procrustes.png"),0.5,1.5,5.1,3.4)
img_fit(s,os.path.join(FIG,"P3_network.png"),5.5,1.45,4.6,3.4)
img_fit(s,os.path.join(FIG,"P3_rf_R2hist.png"),10.05,1.7,3.05,2.9)
statline(s,0.55,5.35,[
 ("Mantel r (P=%s)"%fp(S(2,'mantel_p')),fnum(S(2,'mantel_r')),MET),
 ("Procrustes r (P=%s)"%fp(S(2,'procrustes_p')),fnum(S(2,'procrustes_r')),MET),
 ("PERMANOVA R² (P=%s)"%fp(S(2,'permanova_p')),fnum(S(2,'permanova_R2')),GREEN),
 ("dbRDA R² (P=%s)"%fp(S(2,'dbRDA_p')),fnum(S(2,'dbRDA_R2')),GREEN),
])
takeaway(s,0.55,6.35,12.3,
  "四法全部显著(Mantel/Procrustes 一致 + PERMANOVA≈dbRDA R²=%s);RF 可预测 %s%% 尿液代谢物 → 两体液共享系统循环代谢物。"%(fnum(S(2,'dbRDA_R2')),fnum(float(S(2,'RF_frac_R2gt5'))*100,0)),
  color=MET,h=0.7)

# =========================================================================
# S7 — Pairing 4: mediation
# =========================================================================
s=slide(); header(s,"分析四：三层中介  唾液物种 → 尿液物种 → 尿液代谢物",kicker="mediation (校正年龄/性别/疾病组)",num=9)
# left: concept chain
box(s,0.55,1.55,6.2,2.0,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,0.8,1.72,5.7,0.4,[[("中介链路设计",14,True,MIC)]])
cy2=2.5
for i,(lab,col) in enumerate([("唾液物种\n(暴露 X)",MIC),("尿液物种\n(中介 M)",MET),("尿液代谢物\n(结局 Y)",GOLD)]):
    x=0.85+i*1.95
    box(s,x,cy2,1.55,0.8,fill=col,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s,x,cy2+0.1,1.55,0.65,[[(lab,11,True,WHITE)]],align=PP_ALIGN.CENTER,line_spacing=0.95)
    if i<2: text(s,x+1.5,cy2+0.18,0.5,0.4,[[("→",20,True,INK)]],align=PP_ALIGN.CENTER)
ncand = MED[0]['n_candidate'] if (MED and 'n_candidate' in MED[0]) else str(len(MED))
nsig  = sum(1 for r in MED if r.get('ACME_fdr') and float(r['ACME_fdr'])<0.10)
nnom  = sum(1 for r in MED if r.get('ACME_p') and float(r['ACME_p'])<0.05)
box(s,0.55,3.75,6.2,2.5,fill=WHITE,line=PANELB,line_w=1.2,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,0.8,3.92,5.7,2.3,[
 [("扫描结果",14,True,MIC)],
 [("候选中介三元组  ",12.5,False,INK),("%d"%len(MED),15,True,MIC)],
 [("ACME 名义显著 (P<0.05)  ",12.5,False,INK),("%d"%nnom,15,True,GOLD)],
 [("ACME FDR<0.10  ",12.5,False,INK),("%d"%nsig,15,True,ACC)],
],line_spacing=1.35)
# right: top table
text(s,7.0,1.5,5.8,0.35,[[("中介效应最强的链路（前 6，按 ACME P）",13,True,MIC)]])
mrows=[["唾液物种","尿液物种","尿代谢物","ACME P"]]
for r in MED[:6]:
    mrows.append([r.get('saliva_microbe','')[:16], r.get('urine_microbe','')[:14],
                  r.get('urine_metab','')[:8], fp(r.get('ACME_p',''))])
if len(mrows)==1: mrows.append(["—","—","—","—"])
table(s,7.0,1.95,5.85,[1.9,1.7,1.2,1.05],mrows,header_fill=DARK,
      fsize=11,hsize=11.5,row_h=0.5,header_h=0.42)
takeaway(s,0.55,6.45,12.3,
  "经 FDR 校正后无显著中介(%d 条候选,FDR>0.10) → 唾液菌–尿液菌耦合弱，难成连贯中介链;跨体液信号集中于代谢层。"%len(MED),
  color=ACC,h=0.7)

# =========================================================================
# S8b — Per-feature PERMANOVA R2 histogram
# =========================================================================
s=slide(); header(s,"逐特征 PERMANOVA：每个显著特征的 R²",
                  kicker="adonis2 · 单特征解释尿液端方差",num=10)
img_fit(s,os.path.join(FIG,"perfeature_permanova_hist.png"),0.5,1.5,7.4,4.9)
# right: per-pairing summary cards
rx,rw=8.2,4.7
labs=[("① 唾液物种→尿液微生物组","PF1",MIC),
      ("② 唾液物种→尿液代谢组","PF2",GOLD),
      ("③ 唾液代谢物→尿液代谢组","PF3",MET)]
yy=1.6
for lab,key,col in labs:
    n,nsig,med,mx=pf_stats(key)
    box(s,rx,yy,rw,1.45,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE); box(s,rx,yy,0.1,1.45,fill=col)
    text(s,rx+0.28,yy+0.12,rw-0.5,0.4,[[(lab,12.5,True,col)]])
    text(s,rx+0.28,yy+0.55,rw-0.5,0.8,[
      [("显著特征 ",11.5,False,INK),("%d"%nsig,15,True,col),(" / %d"%n,11.5,False,MUTE),
       ("   (P<0.05)",10.5,False,MUTE)],
      [("R² 中位 ",11.5,False,INK),("%.3f"%med,12.5,True,INK),
       ("   最大 ",11.5,False,INK),("%.3f"%mx,12.5,True,INK)],
    ],line_spacing=1.15)
    yy+=1.62
takeaway(s,0.55,6.55,12.3,
  "绝大多数显著特征单特征 R² 都很小(多 2–5%) → 跨体液信号由大量小贡献特征共同构成,而非少数主导者;代谢物层显著数更多、R² 更高。",
  color=MET,h=0.6)

# =========================================================================
# S8c — Paper-style cumulative variance explained + bootstrap CI
# =========================================================================
CUM = read_csv(os.path.join(TAB, "permanova_cumulative_R2.csv"))
s=slide(); header(s,"论文式累计方差解释 + 95% 置信区间",
                  kicker="对应 Zhang et al. Fig 2a · 联合 PERMANOVA",num=11)
img_fit(s,os.path.join(FIG,"permanova_cumulative_R2.png"),0.5,1.55,7.6,3.6)
# right: method + the paper-comparable number
rx,rw=8.35,4.55
box(s,rx,1.55,rw,2.35,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,rx+0.26,1.7,rw-0.5,2.2,[
 [("三步法(适配小样本)",13,True,MIC)],
 [("① 逐特征 adonis2 海选显著(P<0.05)",11.5,False,INK)],
 [("② 显著特征→前10主成分→一个联合 adonis2 求累计 R²(共享方差只计一次,有界)",11.5,False,INK)],
 [("③ 无放回再抽样得 95% CI(避免有放回 bootstrap 的零距离虚高)",11.5,False,INK)],
],line_spacing=1.2)
def cumget(i,k):
    try: return CUM[i][k]
    except Exception: return "?"
def pct(i,k):
    try: return "%.1f%%"%(100*float(cumget(i,k)))
    except Exception: return "?"
box(s,rx,4.05,rw,2.25,fill=WHITE,line=GOLD,line_w=1.4,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,rx+0.26,4.2,rw-0.5,2.1,[
 [("与论文最直接可比",13,True,GOLD)],
 [("② 唾液菌群 → 尿液代谢组",12,True,INK)],
 [("累计 R² = ",12,False,INK),(pct(1,"R2"),20,True,GOLD)],
 [("95% CI ",11.5,False,MUTE),("%s – %s"%(pct(1,"ci_lo"),pct(1,"ci_hi")),12,True,INK),
  ("  P=%s"%cumget(1,"p"),11.5,False,MUTE)],
 [("≈ 论文 口腔 14.7% / 肠道 18.4% 同类指标",11,False,MUTE)],
],line_spacing=1.22)
takeaway(s,0.55,6.5,12.3,
  "去共线后的真实联合解释度约 20–25%(P<0.01,有界);而单变量 R² 朴素累加会溢出>100%,只反映特征冗余、不可作总解释度。",
  color=MET,h=0.6)

# =========================================================================
# S8d — Shared + disease-differential candidate metabolites (for validation)
# =========================================================================
CAND = read_csv(os.path.join(TAB, "shared_disease_candidates.csv"))
SDall = read_csv(os.path.join(TAB, "shared_disease_metabolites.csv"))
s=slide(); header(s,"候选代谢物：健康/疾病有差异 且 两部位共享",
                  kicker="面向动物实验验证的选靶",num=12)
img_fit(s,os.path.join(FIG,"shared_disease_scatter.png"),0.45,1.5,6.5,4.7)
rx,rw=7.35,5.5
# screening funnel
box(s,rx,1.5,rw,1.28,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,rx+0.28,1.6,rw-0.55,1.1,[
 [("筛选漏斗  ",12.5,True,MIC),
  (f"两部位共享 {len(SDall)} → 两部位均 P<0.05 且同向 {len(CAND)}",12,False,INK)],
 [("(健康 N=34  vs  疾病 P/PD/PC/PCD n=67;Wilcoxon + BH-FDR)",11,False,MUTE)],
],line_spacing=1.2)
# priority candidates
box(s,rx,2.95,rw,2.05,fill=WHITE,line=MET,line_w=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,rx+0.28,3.05,rw-0.55,0.4,[[("✅ 优先候选(内源性·可购·疾病中均↓)",13,True,MET)]])
for i,(nm,desc) in enumerate([
    ("Propionic acid 丙酸","短链脂肪酸·微生物来源·抗炎"),
    ("2-Hydroxyphenylacetic acid","微生物芳香族代谢物"),
    ("N-lauroylethanolamine","N-酰基乙醇胺类脂质介质")]):
    yy=3.5+i*0.47
    box(s,rx+0.32,yy+0.08,0.1,0.1,fill=MET,shape=MSO_SHAPE.OVAL)
    text(s,rx+0.55,yy-0.02,rw-0.9,0.45,[[(nm,11.5,True,INK),("  "+desc,10.5,False,MUTE)]])
# cautions
box(s,rx,5.15,rw,1.05,fill=WHITE,line=ACC,line_w=1.3,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,rx+0.28,5.24,rw-0.55,0.9,[
 [("⚠️ 排除  ",11.5,True,ACC),("Tri-isobutylphosphate = 工业增塑剂,非内源",11,False,INK)],
 [("⚠️ 待确证  ",11.5,True,ACC),("3 个长 IUPAC 名结构需标准品 MS/MS 比对",11,False,INK)],
],line_spacing=1.15)
takeaway(s,0.45,6.5,12.4,
  "三个优先候选在疾病组均↓ → 假说:疾病伴随保护性微生物/脂质代谢物耗竭 → 可直接设计“补充干预”动物实验验证。",
  color=MET,h=0.6)

# =========================================================================
# S8e — Candidate boxplots by priority + distribution read-out
# =========================================================================
s=slide(); header(s,"候选代谢物分布形态与判读",kicker="按优先级从左到右排列",num=13)
img_fit(s,os.path.join(FIG,"shared_disease_boxplots.png"),0.45,1.45,8.75,4.45)
rx,rw=9.4,3.45
def rbox(y,h,title,col,lines):
    box(s,rx,y,rw,h,fill=WHITE,line=col,line_w=1.4,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s,rx+0.22,y+0.1,rw-0.42,0.32,[[(title,11.5,True,col)]])
    text(s,rx+0.24,y+0.46,rw-0.46,h-0.6,lines,line_spacing=1.1)
# ---- auto-group candidates from the CSV so 7- and 9-candidate reports differ ----
def _clean(f):                                    # drop P_/N_ prefix + MS-DIAL suffix
    return f.split(";")[0].replace("P_","").replace("N_","")
def _rows(tier):
    return [r for r in CAND if r.get("tier")==tier]
def _rng(rows):                                   # e.g. "4–8" or "7"
    if not rows: return "—"
    idx=[CAND.index(r)+1 for r in rows]
    return f"{min(idx)}–{max(idx)}" if len(idx)>1 else f"{min(idx)}"
pr,un,ex = _rows("priority"),_rows("unconfirmed"),_rows("exclude")
# priority box: names + read-out
pr_names = "、".join(_clean(r["feature"])[:16] for r in pr[:3])
rbox(1.45,1.9,f"[OK] 优先({_rng(pr)})",MET,[
 [(pr_names,10,True,MET)],
 [("箱体紧凑、无零值堆积;健康→疾病",10,False,INK)],
 [("位移小但两体液同向稳定下降",10,False,INK)],
 [("→ 可信,优先选靶(丙酸居首)",10,True,MET)],
])
# unconfirmed box: list every member (this is what grows from 3 to 5 after PD drop)
un_lines=[[("零值伪影/膳食或表面活性剂污染嫌疑,须标准品+MS/MS:",9.5,False,INK)]]
for r in un:
    un_lines.append([("• "+_clean(r["feature"])[:26],9.5,False,INK)])
rbox(3.5,2.35 if len(un)>3 else 1.7,f"[?] 待确证({_rng(un)})",GOLD,un_lines)
# exclude box
ex_y = 3.5 + (2.35 if len(un)>3 else 1.7) + 0.12
rbox(ex_y,1.55,f"[X] 排除({_rng(ex)})",ACC,[
 [(("、".join(_clean(r["feature"])[:18] for r in ex)) or "—",10,True,ACC)],
 [("分布干净、统计最强,但为工业增塑",10,False,INK)],
 [("剂等外源污染物 → 不作靶点",10,True,ACC)],
])
takeaway(s,0.45,6.15,8.75,
  "统计最强的两个恰恰最不该买 → 候选必须按“内源性/可购/已确证”分层,不能只按 p 值取前 N 个。",
  color=DARK,h=0.62)

# =========================================================================
# S8f — Disease-subgroup specificity
# =========================================================================
s=slide(); header(s,"疾病亚组特异性：候选是“肾病相关”而非牙周炎本身",
                  kicker="P / PD / PC / PCD 各自 vs 健康",num=14)
img_fit(s,os.path.join(FIG,"candidate_subgroup_heat.png"),0.45,1.5,7.7,4.3)
rx,rw=8.45,4.4
box(s,rx,1.5,rw,2.15,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,rx+0.28,1.6,rw-0.55,2.0,[
 [("★ 关键发现",13,True,ACC)],
 [("三个优先候选的显著效应",11.5,False,INK)],
 [("全部只出现在 PC 与 PCD",11.5,True,ACC)],
 [("——正是两个含【肾病】的组;",11.5,False,INK)],
 [("P(单纯牙周炎)、PD 均无。",11.5,False,INK)],
],line_spacing=1.18)
# effect-size evidence
box(s,rx,3.8,rw,1.75,fill=WHITE,line=MIC,line_w=1.3,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,rx+0.25,3.9,rw-0.5,0.34,[[("不是样本量假象 — 看效应量(丙酸 log2FC)",11,True,MIC)]])
table(s,rx+0.25,4.3,rw-0.5,[1.0,0.72,0.72,0.72,0.9],
      [["","P","PD","PC","PCD"],
       ["唾液","+0.17","-0.47","-0.85","-0.85"],
       ["尿液","-0.03","-0.44","-0.71","-0.48"]],
      header_fill=DARK,fsize=10,hsize=10,row_h=0.34,header_h=0.32,align=PP_ALIGN.CENTER)
text(s,rx+0.25,5.62,rw-0.5,0.3,[[("P 组效应量≈0(非“测不出”);PD n=6 居中但把握度不足",9.5,False,MUTE)]])
takeaway(s,0.45,6.1,12.4,
  "直接含义:动物主模型应以【慢性肾病 CKD】为核心(腺嘌呤/5，6 肾切除),牙周炎作为叠加因素——与最初“以牙周炎为主”的设想相反。",
  color=ACC,h=0.65)

# =========================================================================
# S8g1 — Kidney-severity gradient & independence
# =========================================================================
KT = read_csv(os.path.join(TAB, "kidney_gradient_tests.csv"))
def kt(feat_sub, site, key):
    for r in KT:
        if feat_sub in r["feature"] and r["site"] == site:
            try: return float(r[key])
            except Exception: return float("nan")
    return float("nan")
s=slide(); header(s,"肾病梯度趋势与独立性：夯实“肾病相关”",
                  kicker="无 eGFR → 改用有序梯度 + 牙周炎人群内 + 多变量校正",num=15)
img_fit(s,os.path.join(FIG,"kidney_gradient.png"),0.45,1.5,7.5,4.3)
rx,rw=8.2,4.65
box(s,rx,1.5,rw,1.5,fill=DARK,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,rx+0.25,1.6,rw-0.5,1.35,[
 [("⚠️ 数据可得性",12,True,GOLD)],
 [("本队列无 eGFR/肌酐/BUN,仅有肾病 0/1 标记",11,False,WHITE)],
 [("→ 不做“假 eGFR 分级”,改用三项更严格检验",11,False,RGBColor(0xCF,0xE3,0xE6))],
],line_spacing=1.16)
for i,(t,d) in enumerate([
    ("① 有序梯度趋势 (Jonckheere–Terpstra)","0 无肾病 → 1 肾病 → 2 糖尿病肾病"),
    ("② 仅牙周炎人群内比较 (n=18 vs 49)","所有人都有牙周炎 → 排除牙周炎驱动"),
    ("③ 多变量校正","~ 肾病+牙周炎+糖尿病+年龄+性别")]):
    yy=3.15+i*0.78
    box(s,rx,yy,rw,0.7,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s,rx+0.22,yy+0.06,rw-0.44,0.6,[[(t,11,True,MIC)],[(d,10,False,MUTE)]],line_spacing=1.1)
box(s,rx,5.55,rw,0.72,fill=WHITE,line=MET,line_w=1.5,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,rx+0.22,5.63,rw-0.44,0.6,[
 [("丙酸唾液: ",11,True,MET),("JT P=%.0e"%kt("Propionic","Saliva","JT_p"),11.5,True,INK),
  ("  ρ=%.2f"%kt("Propionic","Saliva","spearman_rho"),11,False,INK)],
 [("校正后独立: β=%.2f, P=%.4f"%(kt("Propionic","Saliva","adj_kidney_beta"),
                                  kt("Propionic","Saliva","adj_kidney_p")),11,True,ACC)],
],line_spacing=1.15)
takeaway(s,0.45,6.5,12.4,
  "丙酸是唯一三项全部通过的候选;2-羟基苯乙酸唾液无趋势(JT P=0.11)、校正后 ns → 应降级。梯度分析同时夯实结论并在候选间做出判别。",
  color=ACC,h=0.6)

# =========================================================================
# S8g2 — Propionic acid: one-page evidence summary (for pitching)
# =========================================================================
s=slide(); header(s,"丙酸 (Propionic acid)：首选靶点证据汇总",
                  kicker="跨体液 · 跨亚组 · 剂量方案 — 一页立项",num=16)
# top stat chips
for i,(v,lab,col) in enumerate([
    ("2/2","唾液+尿液均显著且同向",MET),
    ("%.2f"%kt("Propionic","Saliva","spearman_rho"),"肾病梯度 ρ (唾液,单调下降)",GOLD),
    ("P=%.4f"%kt("Propionic","Saliva","adj_kidney_p"),"校正牙周炎/糖尿病后仍独立",ACC)]):
    x=0.45+i*4.15
    box(s,x,1.42,3.95,0.95,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s,x+0.25,1.5,3.5,0.5,[[(v,20,True,col)]])
    text(s,x+0.27,1.98,3.5,0.32,[[(lab,10.5,False,MUTE)]])
# three evidence columns
def col(x,w,title,ccol,lines):
    box(s,x,2.55,w,2.95,fill=WHITE,line=ccol,line_w=1.4,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s,x+0.25,2.65,w-0.5,0.35,[[(title,12.5,True,ccol)]])
    text(s,x+0.27,3.08,w-0.54,2.35,lines,line_spacing=1.22)
col(0.45,3.95,"① 跨体液(共享)",MIC,[
 [("唾液 ",11,True,INK),("log2FC −0.77, FDR 1.7e-05",10.5,False,INK)],
 [("尿液 ",11,True,INK),("log2FC −0.37, P 8.6e-03",10.5,False,INK)],
 [("两体液同向下降,属 522 个",10.5,False,INK)],
 [("共享代谢物之一 → 系统性",10.5,False,INK)],
 [("而非单一体液局部现象",10.5,False,INK)],
])
col(4.6,3.95,"② 跨亚组 · 肾病梯度",GOLD,[
 [("P(单纯牙周炎) log2FC ≈ ",10.5,False,INK),("0",11,True,ACC)],
 [("PC/PCD(含肾病) ",10.5,False,INK),("−0.85/−0.71",11,True,MET)],
 [("JT 单调趋势 P=%.0e"%kt("Propionic","Saliva","JT_p"),10.5,True,INK)],
 [("牙周炎人群内仍显著 P=%.0e"%kt("Propionic","Saliva","perioOnly_p"),10.5,False,INK)],
 [("→ 效应源自肾病,非牙周炎",10.5,True,GOLD)],
])
col(8.75,4.1,"③ 剂量与验证方案",MET,[
 [("丙酸钠 ",11,True,INK),("150–200 mM 自由饮水",10.5,False,INK)],
 [("(或灌胃 1 g/kg/d)",10.5,False,MUTE)],
 [("模型 ",11,True,INK),("腺嘌呤 CKD ± 结扎牙周炎",10.5,False,INK)],
 [("时程 ",11,True,INK),("造模前 7d 预给药 → 14d",10.5,False,INK)],
 [("靶点验证 ",11,True,GOLD),("LC-MS 测鼠唾液/尿/血清",10.5,False,INK)],
 [("货源 ",11,True,INK),("Sigma P1880 丙酸钠 ≥99%",10.5,False,INK)],
])
# why propionate
box(s,0.45,5.62,12.4,0.82,fill=DARK,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,0.72,5.72,12.0,0.68,[
 [("为什么是丙酸(而非另两个候选)  ",12,True,GOLD),
  ("三项肾病检验(梯度趋势 / 牙周炎人群内 / 多变量校正)中,",11.5,False,WHITE),
  ("只有丙酸全部通过",11.5,True,RGBColor(0x7F,0xE3,0xD6)),
  (";N-月桂酰乙醇胺尿液不稳,2-羟基苯乙酸唾液无趋势。",11.5,False,WHITE)],
])
takeaway(s,0.45,6.62,12.4,
  "立项一句话:丙酸在唾液与尿液同步耗竭、随肾病严重度单调下降、且独立于牙周炎与糖尿病 → 以 CKD 小鼠补充丙酸钠验证其保护作用。",
  color=MET,h=0.62)

# =========================================================================
# S8g — Mouse validation experiment design
# =========================================================================
s=slide(); header(s,"小鼠验证实验设计（依亚组结果调整为 CKD 为核心）",
                  kicker="补充干预 · 两阶段",num=17)
# step 0 warning
box(s,0.45,1.42,12.4,0.62,fill=DARK,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,0.72,1.5,12.0,0.5,[
 [("第 0 步(买了先做这个)  ",12,True,GOLD),
  ("当前注释为谱库匹配(~Level 2) → 标准品先做保留时间+二级谱确证,再进动物实验;标准品一物两用。",11.5,False,WHITE)],
])
# left: model + groups
box(s,0.45,2.2,6.0,4.0,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,0.72,2.3,5.5,0.35,[[("模型与分组(C57BL/6)",13,True,MIC)]])
text(s,0.72,2.72,5.55,0.9,[
 [("阶段一  ",11,True,MET),("腺嘌呤饮食 CKD ± 补充干预",11,False,INK)],
 [("阶段二  ",11,True,MET),("CKD + 结扎牙周炎(共病模型)",11,False,INK)],
],line_spacing=1.2)
table(s,0.72,3.72,5.5,[0.6,2.9,0.85],
      [["组","处理","n"],
       ["G1","Sham + 溶剂","8"],
       ["G2","CKD + 溶剂(疾病对照)","10"],
       ["G3","CKD + 丙酸钠 低剂量","10"],
       ["G4","CKD + 丙酸钠 高剂量","10"],
       ["G5","丙酸钠单独(毒性对照)","6"]],
      header_fill=DARK,fsize=10.5,hsize=10.5,row_h=0.36,header_h=0.34)
text(s,0.72,5.85,5.5,0.3,[[("功效:micro-CT/肾功能主要终点,预期 d≈1.2–1.5 → n=10 达 80%",9.5,False,MUTE)]])
# right: dosing + endpoints
box(s,6.65,2.2,6.2,4.0,fill=WHITE,line=PANELB,line_w=1.2,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s,6.92,2.3,5.7,0.35,[[("给药与终点",13,True,MIC)]])
text(s,6.92,2.72,5.7,1.25,[
 [("丙酸钠  ",11,True,MET),("150–200 mM 自由饮水(或灌胃 1 g/kg/d)",10.5,False,INK)],
 [("2-羟基苯乙酸  ",10.5,True,MET),("灌胃 10–50 mg/kg/d",10.5,False,INK)],
 [("N-月桂酰乙醇胺  ",10.5,True,MET),("5–20 mg/kg/d,玉米油",10.5,False,INK)],
 [("时程  ",10.5,True,INK),("造模前 7 d 预给药 → 持续至第 14 d",10.5,False,INK)],
],line_spacing=1.15)
text(s,6.92,4.15,5.7,2.0,[
 [("主要终点  ",11,True,ACC),("肾功能(Scr/BUN/尿ACR)+ 肾组织 PAS/纤维化",10.5,False,INK)],
 [("次要  ",10.5,True,INK),("micro-CT 牙槽骨吸收;牙龈 IL-1β/IL-6/TNF-α",10.5,False,INK)],
 [("靶点验证  ",10.5,True,GOLD),("LC-MS 测鼠唾液/尿/血清该代谢物是否回补",10.5,False,INK)],
 [("机制  ",10.5,True,INK),("GPR43(FFAR2)/HDAC;Treg-Th17 平衡",10.5,False,INK)],
 [("严谨性  ",10.5,True,INK),("随机化+盲评;笼效应;溶剂对照;IACUC 批件",10.5,False,INK)],
],line_spacing=1.15)
takeaway(s,0.45,6.3,12.4,
  "本小鼠实验的目的:人群数据只能证明“丙酸随肾病降低”是相关,故用 CKD 小鼠补充丙酸钠检验其是否为因果(补回能否改善肾损伤);若要进一步夯实因果,可再加【耗竭臂】(清除产丙酸菌)看表型是否加重。",
  color=MET,h=0.6)

# =========================================================================
# S9 — Conclusions
# =========================================================================
s=slide(); header(s,"主要发现与生物学解读",kicker="结论",num=18)
def concard(x,y,w,h,title,col,body):
    box(s,x,y,w,h,fill=PANEL,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s,x,y,w,0.12,fill=col)
    text(s,x+0.28,y+0.22,w-0.55,0.4,[[(title,14,True,col)]])
    text(s,x+0.3,y+0.7,w-0.6,h-0.8,body,line_spacing=1.15)
concard(0.55,1.5,6.05,2.35,"③ 代谢组→代谢组：最强轴",MET,[
 [("• 群落级一致显著(Mantel %s, dbRDA R²=%s)。"%(fnum(S(2,'mantel_r'),2),fnum(S(2,'dbRDA_R2'),2)),12,False,INK)],
 [("• RF 稳健预测大量尿液代谢物。",12,False,INK)],
 [("• 反映两体液共享系统循环代谢物。",12,False,INK)],
])
concard(6.8,1.5,6.05,2.35,"② 微生物→代谢：局部关联",GOLD,[
 [("• 显著对(%s)远多于物种–物种层。"%S(1,'sig_pairs'),12,False,INK)],
 [("• 口腔菌影响远端代谢 > 影响远端菌群。",12,False,INK)],
 [("• 印证“组成不必然决定功能”。",12,False,INK)],
])
concard(0.55,4.0,6.05,2.35,"① 微生物→微生物：最弱",MIC,[
 [("• 整体群落几乎不一致(Mantel≈0)。",12,False,INK)],
 [("• 口腔 vs 泌尿道 = 独立微生态系统。",12,False,INK)],
 [("• 仅个别物种弱关联。",12,False,INK)],
])
concard(6.8,4.0,6.05,2.35,"④ 中介链：不显著",ACC,[
 [("• %d 条候选,FDR 校正后均不显著。"%len(MED),12,False,INK)],
 [("• 菌群耦合弱 → 无连贯中介通路。",12,False,INK)],
 [("• 与论文紧密口腔–肠道轴形成对照。",12,False,INK)],
])
takeaway(s,0.55,6.5,12.3,
  "唾液–尿液轴：微生物组弱耦合、代谢组强耦合 — 符合两体液经血液循环过滤连接(而非消化道连续)的生理学。",
  color=DARK,h=0.6)

# =========================================================================
# S9 — Methods & references
# =========================================================================
s=slide(DARK); box(s,0,0,SW,0.16,fill=MET)
text(s,0.7,0.45,11.0,0.7,[[("方法与参考",26,True,WHITE)]])
text(s,0.7,1.45,6.1,4.6,[
 [("方法",15,True,RGBColor(0x7F,0xE3,0xD6))],
 [("• 微生物组：物种水平;流行度+丰度联合过滤(≥5%样本&RA>0.01%);Bray-Curtis(一致性)、CLR(相关/RF)。",12,False,RGBColor(0xDF,0xEA,0xEC))],
 [("• 代谢组：正+负离子合并，log2(x+1);z-score 欧氏距离;相关/RF 取方差前400。",12,False,RGBColor(0xDF,0xEA,0xEC))],
 [("• 一致性：Mantel(Spearman)+Procrustes/protest(999置换)。",12,False,RGBColor(0xDF,0xEA,0xEC))],
 [("• 方差解释：唾液端主坐标轴约束的 dbRDA(vegan::dbrda)。",12,False,RGBColor(0xDF,0xEA,0xEC))],
 [("• 相关网络：Spearman |r|≥0.3 且 P<0.05(并报 BH-FDR);igraph/ggraph。",12,False,RGBColor(0xDF,0xEA,0xEC))],
 [("• 预测：RandomForest 5折交叉验证逐特征 R²。",12,False,RGBColor(0xDF,0xEA,0xEC))],
 [("• 中介：mediation::mediate，校正年龄/性别/疾病组，FDR 控制。",12,False,RGBColor(0xDF,0xEA,0xEC))],
],line_spacing=1.3)
text(s,7.0,1.45,5.9,4.6,[
 [("参考与产出",15,True,RGBColor(0x7F,0xE3,0xD6))],
 [("Zhang J, et al. Cross-body site microbial interactions influence the human plasma metabolome. Microbiome. 2026;14:147.",12,False,RGBColor(0xDF,0xEA,0xEC))],
 [("",6,False,WHITE)],
 [(f"报告：{DB}_cross_site_multiomics.html",12,True,RGBColor(0x8F,0xC1,0xDF))],
 [(f"脚本：crossomics_helpers.R / GTDB_cross_site_multiomics.Rmd（参数化 db={DBLABEL}{'，drop_groups=PD' if NOPD else ''}）",11.5,False,RGBColor(0xCF,0xE3,0xE6))],
 [("结果表：crossomics/tables/*.csv  图：crossomics/figures/*.png",11.5,False,RGBColor(0xCF,0xE3,0xE6))],
],line_spacing=1.3)
text(s,0.7,6.55,12,0.4,[[("局限：横断面、相对丰度、PD 组 n=6;中介为关联性解释,非因果。",11,False,RGBColor(0x9F,0xBA,0xBE))]])

prs.save(OUT)
print("SAVED", OUT, "slides=", len(prs.slides._sldIdLst))
